"""
build_openclaw_courseware.py  —  CKA-style design (white, blue/green)

Structure:
  Slide 1          : Cover (fresh)
  Slides 2-9       : Admin slides — Digital Attendance, Trainer, Icebreaker,
                     Ground Rules, Schedule, Learning Outcomes, Assessment (fresh)
  Slide 10         : Topic 1 section divider (fresh)
  Slides 11-86     : WSQ v10 PDF pages 12-86 embedded (Topic 1 theory + diagrams)
  Slide 87         : Topic 2 section divider (fresh)
  Slides 88-119    : WSQ v10 PDF pages 88-119 embedded (Topic 2 applications)
  Slide 120        : Topic 3 section divider (fresh)
  Slides 121-149   : WSQ v10 PDF pages 121-149 embedded (lab screenshots)
                     + after each lab header page: fresh COMMANDS + TEST IT slides
  Slides 150+      : WSQ v10 PDF pages 150-163 embedded (closing)
"""

import os
import sys
import shutil
import fitz                               # PyMuPDF
from docx import Document
from docx.shared import Pt, Cm, RGBColor, Inches as DInches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
from pptx import Presentation
from pptx.util import Inches, Pt as PPt
from pptx.dml.color import RGBColor as PRGBColor
from pptx.enum.text import PP_ALIGN

sys.path.insert(0, r"C:\Users\mohan\agents\OpenClaw-repo")
import prodoc

# ── Paths ─────────────────────────────────────────────────────────────────────
REPO       = r"C:\Users\mohan\agents\OpenClaw-repo"
COURSEWARE = os.path.join(REPO, "courseware")
SCREENS    = os.path.join(REPO, "screenshots")
PDF_SRC    = os.path.join(REPO,
    "WSQ  - Learner Guide Slides - Business Transformation"
    " with OpenClaw and NFT - v10.pdf")

os.makedirs(COURSEWARE, exist_ok=True)
os.makedirs(SCREENS,    exist_ok=True)

# ── Course metadata ───────────────────────────────────────────────────────────
COURSE_NAME = "WSQ Business Transformation with OpenClaw"
COURSE_CODE = "TGS-2022015374"
INSTITUTION = "Tertiary Infotech Academy Pte Ltd"
INSTRUCTOR  = "Mohan Pothula"
REFERRAL    = "https://www.hostinger.com?REFERRALCODE=FEGANGCHQ20C"

# ── CKA Design System ─────────────────────────────────────────────────────────
C_WHITE   = PRGBColor(0xFF, 0xFF, 0xFF)
C_BLUE    = PRGBColor(0x25, 0x63, 0xEB)
C_GREEN   = PRGBColor(0x10, 0xB9, 0x81)
C_DARK    = PRGBColor(0x0D, 0x1B, 0x2A)
C_MUTED   = PRGBColor(0x6B, 0x72, 0x80)
C_GHOST   = PRGBColor(0xE5, 0xE7, 0xEB)
C_CODEBG  = PRGBColor(0xF3, 0xF4, 0xF6)
C_LGGREEN = PRGBColor(0xEC, 0xFD, 0xF5)
C_SEPRTR  = PRGBColor(0xD1, 0xD5, 0xDB)
C_BCARD   = PRGBColor(0xDB, 0xEA, 0xFE)

def d(r, g, b): return RGBColor(r, g, b)
D_BLUE  = d(0x25, 0x63, 0xEB)
D_GREEN = d(0x10, 0xB9, 0x81)
D_DARK  = d(0x0D, 0x1B, 0x2A)
D_MUTED = d(0x6B, 0x72, 0x80)

SW = Inches(13.33)
SH = Inches(7.50)

BAR_W  = Inches(0.18)
TOP_H  = Inches(0.04)
BOT_H  = Inches(0.04)
BLUE_H = Inches(4.60)
CX     = BAR_W + Inches(0.20)
CW     = SW - CX - Inches(0.20)
TY     = TOP_H + Inches(0.08)
TITLEY = TOP_H + Inches(0.30)
SEPHY  = TOP_H + Inches(0.88)


# ── Lab data ──────────────────────────────────────────────────────────────────
LABS = [
    ("Lab 1", "OpenClaw Hosting", 30,
     [
      # ── OPTION A: Hostinger VPS ──────────────────────────────────────────
      ("[OPTION A — Hostinger VPS]  1/3 — SSH into your VPS",
       "ssh root@YOUR_VPS_IP\n"
       "# Replace YOUR_VPS_IP with your actual Hostinger VPS IP address"),
      ("[OPTION A — Hostinger VPS]  2/3 — Install OpenClaw",
       "curl -fsSL https://openclaw.ai/install.sh | bash"),
      ("[OPTION A — Hostinger VPS]  3/3 — Install and start the daemon",
       "openclaw daemon install\n"
       "systemctl enable --now openclaw\n"
       "openclaw gateway status   # should show: Gateway: running   Port: 18789"),

      # ── OPTION B: Docker Desktop (local laptop) ───────────────────────────
      ("[OPTION B — Docker Desktop]  1/5 — Open Docker Desktop and confirm it is running",
       "docker version\n"
       "# You should see: Server: Docker Engine - Community\n"
       "# If Docker Desktop is not open, launch it first and wait for 'Engine running'"),
      ("[OPTION B — Docker Desktop]  2/5 — Pull the OpenClaw image from Docker Hub",
       "docker pull openclaw/openclaw:latest\n"
       "# This downloads ~500 MB. Wait for 'Status: Downloaded newer image'"),
      ("[OPTION B — Docker Desktop]  3/5 — Run the OpenClaw container",
       "docker run -d --name openclaw \\\n"
       "  -p 18789:18789 \\\n"
       "  -v openclaw-data:/root/.openclaw \\\n"
       "  openclaw/openclaw:latest"),
      ("[OPTION B — Docker Desktop]  4/5 — Confirm the container is up",
       "docker ps\n"
       "# Expected output (truncated):\n"
       "# CONTAINER ID   IMAGE                      STATUS         PORTS\n"
       "# xxxxxxxxxxxx   openclaw/openclaw:latest   Up X seconds   0.0.0.0:18789->18789/tcp"),
      ("[OPTION B — Docker Desktop]  5/5 — Run the onboarding wizard inside the container",
       "docker exec -it openclaw openclaw onboard"),

      # ── BOTH options ──────────────────────────────────────────────────────
      ("[BOTH OPTIONS]  Open OpenClaw in your browser",
       "http://localhost:18789\n"
       "# Open this URL in Chrome or Firefox on your laptop"),
     ],
     "openclaw gateway status",
     "Gateway: running   Port: 18789",
     "https://killercoda.com/playgrounds/course/kubernetes-playgrounds/ubuntu",
    ),
    ("Lab 2", "OpenClaw Model", 30,
     [("Option A — Groq (free tier, no credit card required)",
       "# Visit console.groq.com → sign up free → API Keys → Create key\n"
       "openclaw model set groq --api-key YOUR_GROQ_KEY --model llama-3.1-70b-versatile"),
      ("Option B — OpenAI OAuth (step 1: launch browser flow)",
       "openclaw model set openai --oauth   # copy the printed URL to your browser"),
      ("Option B — OpenAI OAuth (steps 2-4: authorise and confirm)",
       "# Browser: Sign in to OpenAI → click Allow\n"
       "# Terminal: ✓ OpenAI OAuth token saved."),
      ("Option C — Ollama (local, zero API cost)",
       "ollama launch openclaw\nopenclaw model set ollama --host http://localhost:11434 --model openclaw"),
      ("Option D — Edit config file directly",
       "nano ~/.openclaw/openclaw.json\n# Set: provider, model, apiKey"),
      ("Test the model connection",
       "openclaw model test   # expect: Status: ✓ Connected"),
     ],
     "openclaw model test",
     "Provider: groq   Model: llama-3.1-70b-versatile   Status: ✓ Connected",
     "https://killercoda.com/playgrounds/course/kubernetes-playgrounds/ubuntu",
    ),
    ("Lab 3", "OpenClaw Channel", 30,
     [("Telegram — open BotFather in Telegram app",
       "# Open Telegram → search @BotFather → send /newbot"),
      ("Telegram — enter bot name and username, copy token",
       "# Enter: Alfred Agent  (name)\n# Enter: alfred_agent_bot  (username)\n# Copy the API token shown"),
      ("Telegram — register and start channel",
       "openclaw channel add telegram --token YOUR_BOT_TOKEN\nopenclaw channel start telegram"),
      ("WhatsApp — add channel",
       "openclaw channel add whatsapp"),
      ("WhatsApp — scan QR code on your phone",
       "# QR code appears in terminal\n# WhatsApp → Linked Devices → Link a Device → scan"),
      ("WhatsApp — start channel",
       "openclaw channel start whatsapp"),
     ],
     "openclaw channel list",
     "telegram   running   |   whatsapp   running",
     "https://killercoda.com/playgrounds/course/kubernetes-playgrounds/ubuntu",
    ),
    ("Lab 4", "OpenClaw Tools", 30,
     [("AgentMail — sign up and get API key",
       "# Visit https://agentmail.to → create free account → copy API key"),
      ("AgentMail — add tool",
       "openclaw tools add agentmail --api-key YOUR_AGENTMAIL_KEY"),
      ("Firecrawl — sign up and get fc-... key",
       "# Visit https://www.firecrawl.dev → create account → copy fc-... key"),
      ("Firecrawl — add tool",
       "openclaw tools add firecrawl --api-key fc-YOUR_KEY"),
      ("Agent Browser — get key and add tool",
       "# Visit https://agent-browser.dev → get API key\nopenclaw tools add agent-browser --api-key YOUR_KEY"),
      ("Test tools in chat",
       "# Send: 'Check my email'\n# Send: 'Browse https://news.ycombinator.com'\n# Send: 'Scrape https://openclaw.ai'"),
     ],
     "openclaw tools list",
     "agentmail   enabled   |   firecrawl   enabled   |   agent-browser   enabled",
     "https://killercoda.com/playgrounds/course/kubernetes-playgrounds/ubuntu",
    ),
    ("Lab 5", "OpenClaw Skills", 20,
     [("Browse ClawHub marketplace",
       "# Visit https://clawhub.ai → explore available skills"),
      ("Install web-research skill (VPS / local)",
       "npx clawhub@latest install web-research"),
      ("Install daily-briefing skill (Docker)",
       "docker exec openclaw npx clawhub@latest install daily-briefing"),
      ("List installed skills",
       "openclaw skills list"),
      ("Test skill in Telegram chat",
       "# Send: /web-research What is OpenClaw?"),
      ("Update a skill",
       "npx clawhub@latest update web-research"),
     ],
     "openclaw skills list",
     "web-research   active   |   daily-briefing   active",
     "https://killercoda.com/playgrounds/course/kubernetes-playgrounds/ubuntu",
    ),
    ("Lab 6", "Cron Jobs and Heartbeat", 30,
     [("Create a daily morning summary cron",
       "openclaw cron create --schedule '0 9 * * *' \\\n"
       "  --prompt 'Summarise my email' --channel telegram --name morning"),
      ("Create a weekly AI news cron",
       "openclaw cron create --schedule '0 8 * * 1' \\\n"
       "  --prompt '/web-research Latest AI news' --channel telegram --name weekly"),
      ("List crons and run one immediately",
       "openclaw cron list\nopenclaw cron run morning"),
      ("Configure HEARTBEAT.md",
       "nano ~/.openclaw/HEARTBEAT.md\n# interval: 30m\n# channel: telegram\n# message: OpenClaw is alive"),
      ("Enable heartbeat via CLI",
       "openclaw heartbeat enable --interval 30m --channel telegram"),
      ("Check heartbeat status",
       "openclaw heartbeat status"),
     ],
     "openclaw heartbeat status",
     "Heartbeat: enabled   Interval: 30m   Channel: telegram",
     "https://killercoda.com/playgrounds/course/kubernetes-playgrounds/ubuntu",
    ),
    ("Lab 7", "OpenClaw Security", 30,
     [("Step 1 — Store API keys as environment variables",
       "export GROQ_API_KEY='gsk_...'\necho 'export GROQ_API_KEY=...' >> ~/.bashrc"),
      ("Step 2 — Set Telegram allowlist",
       "openclaw channel config telegram --allowlist-add YOUR_TELEGRAM_USER_ID"),
      ("Step 3 — Set gateway authentication token",
       "openclaw config set gateway.auth-token $(openssl rand -hex 32)"),
      ("Step 4 — Enable execution sandbox",
       "mkdir ~/openclaw-sandbox\nopenclaw tools config exec --sandbox-dir ~/openclaw-sandbox"),
      ("Step 5 — Disable unused tools",
       "openclaw tools disable exec"),
      ("Step 6 — Export and review audit logs",
       "openclaw logs --last 50\nopenclaw logs --export ~/audit.json"),
     ],
     "openclaw config list | grep -E 'gateway|auth'",
     "gateway.auth-token   set   |   gateway.allowlist   enabled",
     "https://killercoda.com/playgrounds/course/kubernetes-playgrounds/ubuntu",
    ),
    ("Lab 8", "OpenClaw Dashboard", 20,
     [("Launch dashboard on local machine",
       "openclaw dashboard   # opens http://localhost:18789"),
      ("VPS — create SSH port-forward tunnel",
       "ssh -L 18789:localhost:18789 root@YOUR_VPS_IP   # keep terminal open"),
      ("VPS — open dashboard in your local browser",
       "# Open: http://localhost:18789 in a NEW browser tab on your LOCAL machine"),
      ("Explore: Overview and Channels panels",
       "# Check: gateway status, uptime, model provider, connected channels"),
      ("Explore: Memory panel",
       "# View MEMORY.md, SOUL.md, AGENTS.md, HEARTBEAT.md"),
      ("Restart gateway from the dashboard",
       "# Dashboard → Gateway → Restart → wait 5 seconds → refresh page"),
     ],
     "openclaw gateway status",
     "Gateway: running   Port: 18789   Dashboard: http://localhost:18789",
     "https://killercoda.com/playgrounds/course/kubernetes-playgrounds/ubuntu",
    ),
    ("Lab 9", "OpenClaw Cost Saving", 20,
     [("Switch to cheapest Claude model (Haiku)",
       "openclaw model set anthropic --model claude-haiku-4-5-20251001"),
      ("Use Claude.ai subscription (fixed monthly fee, unlimited)",
       "openclaw model set anthropic --subscription"),
      ("Use MiniMax free tier (zero API cost)",
       "openclaw model set minimax --api-key KEY --model abab6.5s-chat"),
      ("Use Ollama locally (zero API cost, runs on your machine)",
       "openclaw model set ollama --model openclaw"),
      ("Enable context compaction to reduce token usage",
       "openclaw config set context.compaction enabled\nopenclaw config set context.compaction-threshold 50000"),
      ("Set a monthly spending alert",
       "openclaw usage stats\nopenclaw usage alert --threshold 5.00 --channel telegram"),
     ],
     "openclaw usage stats",
     "Total spend this month: $0.00   |   Compaction: enabled",
     "https://killercoda.com/playgrounds/course/kubernetes-playgrounds/ubuntu",
    ),
    ("Lab 10", "Blockchain Invoice Verification", 30,
     [("Clone the Lab 10 source from GitHub",
       "git clone https://github.com/tertiarycourses/TGS-2022015374-OpenClaw.git\n"
       "cd TGS-2022015374-OpenClaw/labs/lab-10-blockchain-invoice"),
      ("Build the Docker image",
       "docker build -t openclaw/invoice-verify:latest ."),
      ("Run the blockchain verification container",
       "docker run -d -p 5000:5000 --name invoice-verify \\\n"
       "  openclaw/invoice-verify:latest\n"
       "curl http://localhost:5000/health"),
      ("Register a test invoice",
       "curl -X POST http://localhost:5000/register \\\n"
       "  -H 'Content-Type: application/json' \\\n"
       "  -d '{\"invoice_id\":\"INV-2026-001\",\"vendor\":\"Tertiary Infotech\",\n"
       "       \"amount\":1500.00,\"date\":\"2026-07-11\"}'"),
      ("Verify the invoice — same data (should pass)",
       "curl -X POST http://localhost:5000/verify \\\n"
       "  -H 'Content-Type: application/json' \\\n"
       "  -d '{\"invoice_id\":\"INV-2026-001\",\"vendor\":\"Tertiary Infotech\",\n"
       "       \"amount\":1500.00,\"date\":\"2026-07-11\"}'"),
      ("Tamper test — change amount (should fail)",
       "curl -X POST http://localhost:5000/verify \\\n"
       "  -H 'Content-Type: application/json' \\\n"
       "  -d '{\"invoice_id\":\"INV-2026-001\",\"vendor\":\"Tertiary Infotech\",\n"
       "       \"amount\":9999.00,\"date\":\"2026-07-11\"}'\n"
       "# Expected: \"verified\": false, \"tampered\": true"),
     ],
     "curl http://localhost:5000/health",
     '{"invoices": 1, "status": "ok"}',
     "https://github.com/tertiarycourses/TGS-2022015374-OpenClaw/tree/main/labs/lab-10-blockchain-invoice",
    ),
]

TOTAL_MINS = sum(l[2] for l in LABS)


# ── PPTX helpers ──────────────────────────────────────────────────────────────
def blank_slide(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = C_WHITE
    return sl


def add_rect(sl, x, y, w, h, colour):
    s = sl.shapes.add_shape(1, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = colour
    s.line.fill.background()
    return s


def add_oval(sl, x, y, w, h, colour):
    s = sl.shapes.add_shape(9, x, y, w, h)  # 9 = OVAL
    s.fill.solid()
    s.fill.fore_color.rgb = colour
    s.line.fill.background()
    return s


def add_tb(sl, text, x, y, w, h, fname, fsize, bold, colour,
           align=PP_ALIGN.LEFT, italic=False):
    txb = sl.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name   = fname
    run.font.size   = PPt(fsize)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = colour
    return txb


def cka_chrome(sl):
    """Thin blue top strip + thin green bottom strip + left bar (blue top, green bottom)."""
    add_rect(sl, 0, 0, SW, TOP_H, C_BLUE)
    add_rect(sl, 0, SH - BOT_H, SW, BOT_H, C_GREEN)
    add_rect(sl, 0, TOP_H, BAR_W, BLUE_H, C_BLUE)
    add_rect(sl, 0, TOP_H + BLUE_H, BAR_W, SH - TOP_H - BOT_H - BLUE_H, C_GREEN)


def footer(sl, pg):
    add_tb(sl, f"{COURSE_NAME}  ·  {COURSE_CODE}",
           Inches(0.30), SH - Inches(0.32), Inches(7.5), Inches(0.25),
           "Calibri", 9, False, C_MUTED, PP_ALIGN.LEFT)
    add_tb(sl, "© 2026 Tertiary Infotech Academy Pte Ltd",
           Inches(7.80), SH - Inches(0.32), Inches(4.0), Inches(0.25),
           "Calibri", 9, False, C_MUTED, PP_ALIGN.RIGHT)
    add_tb(sl, str(pg),
           Inches(12.90), SH - Inches(0.32), Inches(0.35), Inches(0.25),
           "Calibri", 9, False, C_MUTED, PP_ALIGN.RIGHT)


def slide_from_image(prs, img_path):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    sl.shapes.add_picture(img_path, 0, 0, SW, SH)
    # White bar covers the original PDF footer; thin green strip for CKA chrome consistency
    add_rect(sl, 0, SH - Inches(0.30), SW, Inches(0.26), C_WHITE)
    add_rect(sl, 0, SH - BOT_H, SW, BOT_H, C_GREEN)
    return sl


# ── Slide builders ────────────────────────────────────────────────────────────
def cka_cover(prs, pg):
    sl = blank_slide(prs)
    add_rect(sl, 0, 0, SW, TOP_H, C_BLUE)
    add_rect(sl, 0, SH - BOT_H, SW, BOT_H, C_GREEN)
    add_tb(sl, "COURSE SLIDES",
           Inches(0.75), Inches(2.35), Inches(10), Inches(0.28),
           "Calibri", 12, True, C_BLUE)
    add_tb(sl, COURSE_NAME,
           Inches(0.75), Inches(2.75), Inches(11.5), Inches(1.60),
           "Calibri", 40, True, C_DARK)
    add_rect(sl, Inches(0.75), Inches(4.30), Inches(2.90), Inches(0.05), C_GREEN)
    add_tb(sl, f"WSQ Course Code: {COURSE_CODE}",
           Inches(0.75), Inches(4.50), Inches(9), Inches(0.32),
           "Calibri", 15, False, C_DARK)
    add_tb(sl, f"Conducted by {INSTITUTION}  ·  UEN 201200696W",
           Inches(0.75), Inches(4.88), Inches(9), Inches(0.32),
           "Calibri", 15, False, C_DARK)
    add_tb(sl, "© 2026 Tertiary Infotech Academy Pte Ltd  ·  www.tertiarycourses.com.sg",
           Inches(0.75), Inches(7.05), Inches(11.5), Inches(0.25),
           "Calibri", 9, False, C_MUTED)


def cka_section_divider(prs, section_label, title, pg):
    sl = blank_slide(prs)
    cka_chrome(sl)
    add_tb(sl, section_label, CX, TY, CW, Inches(0.24), "Calibri", 11, True, C_BLUE)
    add_tb(sl, title, CX, TITLEY, CW, Inches(0.60), "Calibri", 36, True, C_DARK)
    footer(sl, pg)


def cka_topic_divider(prs, label, title, subtitle, ghost, pg):
    sl = blank_slide(prs)
    cka_chrome(sl)
    add_tb(sl, ghost,
           Inches(9.50), Inches(0.10), Inches(3.60), Inches(2.10),
           "Calibri", 130, True, C_GHOST, PP_ALIGN.RIGHT)
    add_tb(sl, label, CX, Inches(2.40), CW, Inches(0.28), "Calibri", 12, True, C_BLUE)
    add_tb(sl, title, CX, Inches(2.78), CW, Inches(1.60), "Calibri", 44, True, C_DARK)
    add_tb(sl, subtitle, CX, Inches(4.35), CW, Inches(0.35), "Calibri", 16, False, C_MUTED)
    footer(sl, pg)


def cka_bullet_slide(prs, section_label, title, bullets, pg):
    sl = blank_slide(prs)
    cka_chrome(sl)
    add_tb(sl, section_label, CX, TY, CW, Inches(0.24), "Calibri", 11, True, C_BLUE)
    add_tb(sl, title, CX, TITLEY, CW, Inches(0.52), "Calibri", 30, True, C_DARK)
    add_rect(sl, CX, SEPHY, CW, Inches(0.02), C_SEPRTR)
    y = SEPHY + Inches(0.12)
    for b in bullets:
        add_tb(sl, b, CX, y, CW, Inches(0.40), "Calibri", 13, False, C_DARK)
        y += Inches(0.43)
    footer(sl, pg)


def cka_cards_slide(prs, section_label, title, intro, cards, pg):
    sl = blank_slide(prs)
    cka_chrome(sl)
    add_tb(sl, section_label, CX, TY, CW, Inches(0.24), "Calibri", 11, True, C_BLUE)
    add_tb(sl, title, CX, TITLEY, CW, Inches(0.52), "Calibri", 30, True, C_DARK)
    sep_y = SEPHY + (Inches(0.22) if intro else 0)
    if intro:
        add_tb(sl, intro, CX, SEPHY - Inches(0.04), CW, Inches(0.26),
               "Calibri", 12, False, C_MUTED)
    add_rect(sl, CX, sep_y, CW, Inches(0.02), C_SEPRTR)
    card_w = Inches(3.80)
    card_h = Inches(1.65)
    gap_x  = Inches(0.28)
    gap_y  = Inches(0.20)
    row_y  = sep_y + Inches(0.22)
    for i, (ct, cb) in enumerate(cards[:6]):
        col = i % 3
        row = i // 3
        cx  = CX + col * (card_w + gap_x)
        cy  = row_y + row * (card_h + gap_y)
        add_rect(sl, cx, cy, card_w, card_h, C_CODEBG)
        add_rect(sl, cx + Inches(0.15), cy + Inches(0.15),
                 Inches(0.38), Inches(0.38), C_BLUE)
        add_tb(sl, ct,
               cx + Inches(0.65), cy + Inches(0.12), card_w - Inches(0.75), Inches(0.32),
               "Calibri", 12, True, C_DARK)
        add_tb(sl, cb,
               cx + Inches(0.15), cy + Inches(0.60), card_w - Inches(0.25), Inches(0.95),
               "Calibri", 11, False, C_MUTED)
    footer(sl, pg)


def _code_block(sl, label, command, x, y, w):
    """Draw step label (blue) + code block (gray bg, green left border)."""
    add_tb(sl, label, x, y, w, Inches(0.24), "Calibri", 10, True, C_BLUE)
    y += Inches(0.25)
    lines   = command.split("\n")
    block_h = Inches(0.20) + Inches(0.22) * max(len(lines), 1)
    add_rect(sl, x, y, Inches(0.07), block_h, C_GREEN)
    add_rect(sl, x + Inches(0.07), y, w - Inches(0.07), block_h, C_CODEBG)
    add_tb(sl, command,
           x + Inches(0.14), y + Inches(0.04),
           w - Inches(0.22), block_h - Inches(0.08),
           "Courier New", 9, False, C_DARK)
    return y + block_h + Inches(0.12)


def cka_lab_commands(prs, lab_num, lab_title, steps, pg, cont=False):
    sl = blank_slide(prs)
    cka_chrome(sl)
    add_tb(sl, f"{lab_num}  ·  COMMANDS",
           CX, TY, CW, Inches(0.24), "Calibri", 11, True, C_GREEN)
    add_tb(sl, lab_title + (" (cont.)" if cont else ""),
           CX, TITLEY, CW, Inches(0.50), "Calibri", 26, True, C_DARK)
    add_rect(sl, CX, SEPHY, CW, Inches(0.02), C_SEPRTR)
    y = SEPHY + Inches(0.14)
    for i, (step_label, command) in enumerate(steps, 1):
        y = _code_block(sl, f"Step {i}  —  {step_label}", command, CX, y, CW)
        if y > Inches(6.65):
            break
    footer(sl, pg)


def cka_lab_testit(prs, lab_num, lab_title, test_cmd, expected, url, pg):
    sl = blank_slide(prs)
    cka_chrome(sl)
    add_tb(sl, f"{lab_num}  ·  TEST IT",
           CX, TY, CW, Inches(0.24), "Calibri", 11, True, C_GREEN)
    add_tb(sl, lab_title, CX, TITLEY, CW, Inches(0.50), "Calibri", 26, True, C_DARK)
    add_rect(sl, CX, SEPHY, CW, Inches(0.02), C_SEPRTR)
    card_y = SEPHY + Inches(0.20)
    card_h = Inches(2.90)
    add_rect(sl, CX, card_y, CW, card_h, C_LGGREEN)
    add_tb(sl, "✓  Test it",
           CX + Inches(0.25), card_y + Inches(0.20), Inches(5), Inches(0.40),
           "Calibri", 18, True, C_GREEN)
    add_tb(sl, test_cmd,
           CX + Inches(0.25), card_y + Inches(0.70), CW - Inches(0.50), Inches(0.32),
           "Courier New", 12, False, C_DARK)
    add_tb(sl, f"Expected:  {expected}",
           CX + Inches(0.25), card_y + Inches(1.12), CW - Inches(0.50), Inches(0.32),
           "Calibri", 12, False, C_MUTED)
    if url:
        add_tb(sl, "Killercoda:", Inches(0.30), Inches(6.56), Inches(1.20), Inches(0.28),
               "Calibri", 11, True, C_GREEN)
        add_tb(sl, url, Inches(1.58), Inches(6.56), Inches(11.0), Inches(0.28),
               "Calibri", 11, False, C_GREEN)
    footer(sl, pg)


# ── Admin slide helpers ───────────────────────────────────────────────────────
def slide_digital_attendance(prs, pg):
    sl = blank_slide(prs)
    cka_chrome(sl)
    add_tb(sl, "TRAQOM  ·  SSG DIGITAL ATTENDANCE",
           CX, TY, CW, Inches(0.24), "Calibri", 11, True, C_BLUE)
    add_tb(sl, "Digital Attendance (Mandatory)",
           CX, TITLEY, CW, Inches(0.52), "Calibri", 30, True, C_DARK)
    add_rect(sl, CX, SEPHY, CW, Inches(0.02), C_SEPRTR)
    bullets = [
        "•  It is mandatory to take the AM, PM and Assessment digital attendance for WSQ-funded courses.",
        "•  The trainer displays the digital attendance QR code from the SSG portal.",
        "•  Scan the QR code with your phone camera and submit your attendance.",
        "•  A minimum of 75% attendance is required for assessment and funding.",
    ]
    y = SEPHY + Inches(0.14)
    for b in bullets:
        add_tb(sl, b, CX, y, Inches(7.5), Inches(0.42), "Calibri", 13, False, C_DARK)
        y += Inches(0.46)
    add_rect(sl, Inches(9.40), SEPHY + Inches(0.14), Inches(3.60), Inches(2.20), C_BCARD)
    add_tb(sl, "Minimum 75%\nattendance required",
           Inches(9.50), SEPHY + Inches(0.55), Inches(3.40), Inches(1.20),
           "Calibri", 18, True, C_BLUE, PP_ALIGN.CENTER)
    footer(sl, pg)


def slide_about_trainer(prs, pg):
    sl = blank_slide(prs)
    cka_chrome(sl)
    add_tb(sl, "YOUR TRAINER", CX, TY, CW, Inches(0.24), "Calibri", 11, True, C_BLUE)
    add_tb(sl, "About the Trainer", CX, TITLEY, CW, Inches(0.52), "Calibri", 30, True, C_DARK)
    add_rect(sl, CX, SEPHY, CW, Inches(0.02), C_SEPRTR)

    # ── Left dark card ──
    CARD_X = CX
    CARD_Y = SEPHY + Inches(0.20)
    CARD_W = Inches(3.20)
    CARD_H = Inches(4.40)
    add_rect(sl, CARD_X, CARD_Y, CARD_W, CARD_H, C_DARK)

    # Circular avatar (blue, centred in card)
    AV_D = Inches(1.20)
    AV_X = CARD_X + (CARD_W - AV_D) / 2
    AV_Y = CARD_Y + Inches(0.30)
    add_oval(sl, AV_X, AV_Y, AV_D, AV_D, C_BLUE)
    add_tb(sl, "👤", AV_X, AV_Y + Inches(0.10), AV_D, AV_D - Inches(0.10),
           "Segoe UI Emoji", 28, False, C_WHITE, PP_ALIGN.CENTER)

    # Name
    add_tb(sl, INSTRUCTOR,
           CARD_X + Inches(0.10), AV_Y + AV_D + Inches(0.20),
           CARD_W - Inches(0.20), Inches(0.45),
           "Calibri", 16, True, C_WHITE, PP_ALIGN.CENTER)
    # Subtitle
    add_tb(sl, "Agentic AI · Generative AI · Cloud Trainer",
           CARD_X + Inches(0.10), AV_Y + AV_D + Inches(0.70),
           CARD_W - Inches(0.20), Inches(0.40),
           "Calibri", 10, False, C_BLUE, PP_ALIGN.CENTER)

    # ── Right detail rows ──
    ICON_D  = Inches(0.55)
    ROW_X   = CARD_X + CARD_W + Inches(0.40)
    TEXT_X  = ROW_X + ICON_D + Inches(0.20)
    TEXT_W  = SW - TEXT_X - Inches(0.20)
    ROW_Y0  = CARD_Y + Inches(0.20)
    ROW_GAP = Inches(1.00)

    details = [
        ("👤", "Role",           "Trainer at Tertiary Infotech Academy Pte Ltd"),
        ("🎓", "Qualifications", "PhD  ·  AI  ·  Cloud  ·  Docker  ·  Kubernetes"),
        ("💻", "Expertise",      "Agentic AI  ·  RAG  ·  Generative AI  ·  OpenClaw  ·  LLM  ·  DevOps"),
        ("💼", "Experience",     "20+ years across industry and academia"),
    ]
    for i, (icon, label, val) in enumerate(details):
        iy = ROW_Y0 + i * ROW_GAP
        add_oval(sl, ROW_X, iy, ICON_D, ICON_D, C_BLUE)
        add_tb(sl, icon, ROW_X, iy + Inches(0.08), ICON_D, ICON_D - Inches(0.08),
               "Segoe UI Emoji", 16, False, C_WHITE, PP_ALIGN.CENTER)
        add_tb(sl, label, TEXT_X, iy, TEXT_W, Inches(0.32),
               "Calibri", 13, True, C_DARK)
        add_tb(sl, val, TEXT_X, iy + Inches(0.32), TEXT_W, Inches(0.55),
               "Calibri", 12, False, C_MUTED)

    footer(sl, pg)


def slide_icebreaker(prs, pg):
    cka_cards_slide(prs, "ICE-BREAKER", "Let's Know Each Other",
                    "Take a minute to introduce yourself to the class:", [
                        ("Your role",       "Your name and organisation / role."),
                        ("Your experience", "Any experience with AI tools or automation."),
                        ("Your goal",       "What you hope to build with OpenClaw."),
                    ], pg)


def slide_ground_rules(prs, pg):
    cka_cards_slide(prs, "HOUSE RULES", "Ground Rules", None, [
        ("Phones on silent",    "Set your phone to silent mode."),
        ("Participate actively","No question is too small — ask freely."),
        ("Mutual respect",      "One conversation at a time."),
        ("Be punctual",         "Return from breaks on time."),
        ("Step out quietly",    "For calls or short breaks."),
        ("75% attendance",      "Required for funding eligibility."),
    ], pg)


def slide_schedule(prs, pg):
    sl = blank_slide(prs)
    cka_chrome(sl)
    add_tb(sl, "SCHEDULE", CX, TY, CW, Inches(0.24), "Calibri", 11, True, C_BLUE)
    add_tb(sl, "Lesson Plan — 1 Day, 7½ hours  (09:30 – 17:00)",
           CX, TITLEY, CW, Inches(0.52), "Calibri", 30, True, C_DARK)
    add_rect(sl, CX, SEPHY, CW, Inches(0.02), C_SEPRTR)
    slots_l = [
        ("09:30", "Welcome, Digital Attendance & Overview"),
        ("09:50", "Topic 1: Overview of OpenClaw  [30 min]"),
        ("10:20", "Topic 3: Tools, Skills & Channels  [30 min]"),
        ("10:50", "Tea Break"),
        ("11:05", "Topic 4: Security, Cost Saving & Best Practices  [20 min]"),
        ("11:25", "Topic 2: Applications + Blockchain Demo  [20 min]"),
        ("11:45", "Lab 1: OpenClaw Hosting  [30 min]"),
        ("12:15", "Lab 2: OpenClaw Model (Groq / OpenAI / Ollama)  [30 min]"),
        ("12:45", "Lab 3: OpenClaw Channel  [20 min]"),
        ("13:05", "Lunch  (1 hour)"),
    ]
    slots_r = [
        ("14:05", "Lab 4: OpenClaw Tools  [30 min]"),
        ("14:35", "Lab 5: OpenClaw Skills  [20 min]"),
        ("14:55", "Lab 6: Cron Jobs and Heartbeat  [30 min]"),
        ("15:25", "Lab 7: OpenClaw Security  [30 min]"),
        ("15:55", "Lab 8: OpenClaw Dashboard  [20 min]"),
        ("16:15", "Lab 9: Cost Saving  [20 min]"),
        ("16:35", "TRAQOM Survey, Certificate & Assessment"),
        ("17:00", "End of Course"),
    ]
    y_l = SEPHY + Inches(0.16)
    for time, act in slots_l:
        add_tb(sl, time, CX, y_l, Inches(0.75), Inches(0.30), "Calibri", 11, True, C_BLUE)
        add_tb(sl, act, CX + Inches(0.82), y_l, Inches(5.70), Inches(0.30), "Calibri", 11, False, C_DARK)
        y_l += Inches(0.35)
    y_r = SEPHY + Inches(0.16)
    rx = CX + Inches(7.00)
    for time, act in slots_r:
        add_tb(sl, time, rx, y_r, Inches(0.75), Inches(0.30), "Calibri", 11, True, C_BLUE)
        add_tb(sl, act, rx + Inches(0.82), y_r, Inches(5.30), Inches(0.30), "Calibri", 11, False, C_DARK)
        y_r += Inches(0.35)
    add_rect(sl, CX, Inches(6.80), CW, Inches(0.32), C_CODEBG)
    add_tb(sl, "Daily timing:  9:30am – 5:00pm  ·  1-hour lunch  ·  15-min tea break",
           CX + Inches(0.15), Inches(6.82), CW - Inches(0.30), Inches(0.28),
           "Calibri", 10, False, C_MUTED)
    footer(sl, pg)


def slide_learning_outcomes(prs, pg):
    cka_bullet_slide(prs, "LEARNING OUTCOMES", "What You Will Learn", [
        "•  LO1: Deploy OpenClaw on a Hostinger VPS or Docker Desktop local machine",
        "•  LO2: Connect LLM providers — Groq (free), OpenAI OAuth, Ollama, Default config",
        "•  LO3: Set up Telegram (BotFather) and WhatsApp (QR pairing) channels",
        "•  LO4: Integrate tools — AgentMail, Agent Browser, Firecrawl",
        "•  LO5: Install and invoke skills from the ClawHub marketplace",
        "•  LO6: Schedule tasks with cron jobs and monitor uptime with heartbeat",
        "•  LO7: Apply 10-step security hardening for the OpenClaw gateway",
        "•  LO8: Access the OpenClaw web dashboard (local or via SSH tunnel)",
        "•  LO9: Optimise costs via model selection and context compaction",
        "•  LO10: Implement blockchain invoice verification using Docker + OpenClaw tools",
    ], pg)


def slide_assessment(prs, pg):
    sl = blank_slide(prs)
    cka_chrome(sl)
    add_tb(sl, "FINAL ASSESSMENT", CX, TY, CW, Inches(0.24), "Calibri", 11, True, C_BLUE)
    add_tb(sl, "Assessment", CX, TITLEY, CW, Inches(0.52), "Calibri", 30, True, C_DARK)
    add_rect(sl, CX, SEPHY, CW, Inches(0.02), C_SEPRTR)
    # Left card
    add_rect(sl, CX, SEPHY + Inches(0.18), Inches(5.80), Inches(3.40), C_CODEBG)
    add_tb(sl, "Final Assessment",
           CX + Inches(0.20), SEPHY + Inches(0.30), Inches(5.40), Inches(0.35),
           "Calibri", 14, True, C_DARK)
    items = ["Written / MCQ assessment on the LMS",
             "Open book — slides and Learner Guide allowed",
             "Must be attempted for SSG funding",
             "Covers all 10 labs and OpenClaw concepts"]
    ay = SEPHY + Inches(0.72)
    for item in items:
        add_tb(sl, f"•  {item}", CX + Inches(0.20), ay, Inches(5.40), Inches(0.36),
               "Calibri", 12, False, C_DARK)
        ay += Inches(0.38)
    # Right blue card
    add_rect(sl, CX + Inches(6.30), SEPHY + Inches(0.18), Inches(6.20), Inches(3.40), C_BLUE)
    add_tb(sl, "Funding & Competency",
           CX + Inches(6.55), SEPHY + Inches(0.30), Inches(5.80), Inches(0.40),
           "Calibri", 16, True, C_WHITE)
    add_tb(sl, "Minimum 75% attendance",
           CX + Inches(6.55), SEPHY + Inches(0.90), Inches(5.80), Inches(0.30),
           "Calibri", 13, True, C_WHITE)
    add_tb(sl, "Based on SSG digital attendance records.",
           CX + Inches(6.55), SEPHY + Inches(1.22), Inches(5.80), Inches(0.28),
           "Calibri", 11, False, C_WHITE)
    add_tb(sl, "Assessed as Competent",
           CX + Inches(6.55), SEPHY + Inches(1.82), Inches(5.80), Inches(0.30),
           "Calibri", 13, True, C_WHITE)
    add_tb(sl, "Pass the written / MCQ assessment.",
           CX + Inches(6.55), SEPHY + Inches(2.14), Inches(5.80), Inches(0.28),
           "Calibri", 11, False, C_WHITE)
    footer(sl, pg)


# ── Phase 1: Render PDF pages ────────────────────────────────────────────────
def render_pdf_pages():
    print("  Rendering WSQ v10 PDF pages …")
    doc = fitz.open(PDF_SRC)
    mat = fitz.Matrix(2.5, 2.5)
    paths = {}
    for i in range(doc.page_count):
        out = os.path.join(SCREENS, f"slide_{i+1:03d}.png")
        if not os.path.exists(out):
            pix = doc[i].get_pixmap(matrix=mat)
            pix.save(out)
        paths[i + 1] = out
    print(f"  {len(paths)} pages ready in {SCREENS}")
    return paths


# ── Phase 2: Build PPTX ──────────────────────────────────────────────────────
LAB_PAGES = {121: 0, 125: 1, 129: 2, 133: 3, 136: 4,
             137: 5, 138: 6, 140: 7, 143: 8}

# PDF pages to skip entirely (no image embedded; LAB_PAGES trigger still fires if applicable)
SKIP_PAGES = {21,                                               # slide 20 (not required)
              24, 25, 37, 38, 80, 81,                          # removed theory slides
              92, 93, 94, 95, 96, 97, 98, 99, 100, 101, 102,  # removed slides 91-101
              125,                                              # slide 110 MiniMax lab header → replaced by fresh Lab 2
              145, 146, 147, 148, 149}                         # Google Integration removed


def build_ppt(page_images):
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH
    pg = 1

    # Admin (fresh CKA)
    cka_cover(prs, pg);                pg += 1
    cka_section_divider(prs, "COURSE ADMINISTRATION", "Welcome & Housekeeping", pg); pg += 1
    slide_digital_attendance(prs, pg); pg += 1
    slide_about_trainer(prs, pg);      pg += 1
    slide_icebreaker(prs, pg);         pg += 1
    slide_ground_rules(prs, pg);       pg += 1
    slide_schedule(prs, pg);           pg += 1
    slide_learning_outcomes(prs, pg);  pg += 1
    slide_assessment(prs, pg);         pg += 1

    # Topic 1
    cka_topic_divider(prs, "TOPIC 1", "Overview of OpenClaw",
                      "AI Agent evolution · Context · HEARTBEAT · Gateway · Memory",
                      "01", pg); pg += 1
    for pnum in range(12, 87):
        if pnum in SKIP_PAGES:
            continue
        slide_from_image(prs, page_images[pnum]); pg += 1

    # Topic 2
    cka_topic_divider(prs, "TOPIC 2", "OpenClaw Applications",
                      "AgentMail · Agent Browser · Firecrawl · Mission Control · NFT",
                      "02", pg); pg += 1
    for pnum in range(88, 120):
        if pnum in SKIP_PAGES:
            continue
        slide_from_image(prs, page_images[pnum]); pg += 1

    # Topic 3 / Labs
    cka_topic_divider(prs, "TOPIC 3", "OpenClaw Setup and Configurations",
                      "10 hands-on labs · Hostinger VPS and Docker Desktop environments",
                      "03", pg); pg += 1

    labs_inserted = set()
    for pnum in range(121, 150):
        if pnum not in page_images:
            continue
        if pnum not in SKIP_PAGES:
            slide_from_image(prs, page_images[pnum]); pg += 1
        if pnum in LAB_PAGES:
            idx = LAB_PAGES[pnum]
            if idx not in labs_inserted:
                num, title, mins, steps, test_cmd, expected, url = LABS[idx]
                groups = [steps[i:i+4] for i in range(0, len(steps), 4)]
                for gi, grp in enumerate(groups):
                    cka_lab_commands(prs, num, title, grp, pg, cont=(gi > 0))
                    pg += 1
                cka_lab_testit(prs, num, title, test_cmd, expected, url, pg)
                pg += 1
                labs_inserted.add(idx)

    # Lab 10 — Blockchain Invoice Verification (no PDF trigger page, injected directly)
    num, title, mins, steps, test_cmd, expected, url = LABS[9]
    groups = [steps[i:i+4] for i in range(0, len(steps), 4)]
    for gi, grp in enumerate(groups):
        cka_lab_commands(prs, num, title, grp, pg, cont=(gi > 0))
        pg += 1
    cka_lab_testit(prs, num, title, test_cmd, expected, url, pg)
    pg += 1

    # Closing pages from v10 PDF
    for pnum in range(150, 164):
        if pnum in page_images:
            slide_from_image(prs, page_images[pnum]); pg += 1

    out = os.path.join(COURSEWARE, "WSQ-Business-Transformation-with-OpenClaw.pptx")
    prs.save(out)
    sz = os.path.getsize(out) // 1024
    print(f"  Saved: {os.path.basename(out)}  ({pg-1} slides, {sz:,} KB)")
    return out


# ── LG helpers ───────────────────────────────────────────────────────────────
_LG_BRAND = RGBColor(0x1F, 0x6F, 0xEB)
_LG_DARK  = RGBColor(0x11, 0x18, 0x27)
_LG_GREY  = RGBColor(0x55, 0x5B, 0x66)

def _lg_shade_cell(cell, hexc):
    tcPr = cell._tc.get_or_add_tcPr(); shd = OxmlElement("w:shd")
    shd.set(qn("w:val"), "clear"); shd.set(qn("w:color"), "auto"); shd.set(qn("w:fill"), hexc); tcPr.append(shd)

def _lg_shade_para(pr, hexc="F3F5F8"):
    ppr = pr._p.get_or_add_pPr(); shd = OxmlElement("w:shd")
    shd.set(qn("w:val"), "clear"); shd.set(qn("w:color"), "auto"); shd.set(qn("w:fill"), hexc); ppr.append(shd)

def _lg_runs(paragraph, text):
    import re
    for part in re.split(r"(\*\*[^*]+\*\*|`[^`]+`)", text):
        if not part: continue
        if part.startswith("**") and part.endswith("**"):
            r = paragraph.add_run(part[2:-2]); r.bold = True
        elif part.startswith("`") and part.endswith("`"):
            r = paragraph.add_run(part[1:-1]); r.font.name = "Consolas"; r.font.color.rgb = RGBColor(0xC7, 0x25, 0x4E)
        else:
            paragraph.add_run(part)


# ── LG ───────────────────────────────────────────────────────────────────────
def build_lg():
    ASSETS = os.path.join(REPO, "courseware", "assets")
    VERSIONS = [("1.0", "July 2026",
                 "First version — step-by-step guide to all 10 OpenClaw labs: "
                 "hosting, LLM models, channels, tools, skills, cron jobs, security, "
                 "dashboard, cost saving and blockchain invoice verification",
                 INSTITUTION)]

    doc = Document()
    doc.styles["Normal"].font.name = "Arial"; doc.styles["Normal"].font.size = Pt(11)
    prodoc.style_headings(doc)
    prodoc.add_cover_page(doc, "Learner Guide", COURSE_NAME, "1.0",
                          org_logo=os.path.join(ASSETS, "tertiary-infotech-logo.png"))
    prodoc.add_version_control(doc, VERSIONS)
    prodoc.add_toc(doc, levels="1-2")

    def h1(t): doc.add_paragraph(style="Heading 1").add_run(t)
    def h2(t): doc.add_paragraph(style="Heading 2").add_run(t)

    def p(text):
        _lg_runs(doc.add_paragraph(), text)

    def bullets(items):
        for s in items:
            pr = doc.add_paragraph(style="List Bullet"); pr.paragraph_format.space_after = Pt(2)
            _lg_runs(pr, s)

    def code(text):
        pr = doc.add_paragraph(); _lg_shade_para(pr)
        r = pr.add_run(text); r.font.name = "Consolas"; r.font.size = Pt(9)

    def note(text):
        pr = doc.add_paragraph(); _lg_shade_para(pr, "FFF4E5")
        rr = pr.add_run("Note:  "); rr.bold = True; rr.font.color.rgb = RGBColor(0xB5, 0x6A, 0x00)
        _lg_runs(pr, text)

    def test_it(text):
        pr = doc.add_paragraph(); _lg_shade_para(pr, "E8F7EE")
        rr = pr.add_run("Test it:  "); rr.bold = True; rr.font.color.rgb = RGBColor(0x12, 0x7A, 0x3E)
        _lg_runs(pr, text)

    def rule():
        pr = doc.add_paragraph()
        ppr = pr._p.get_or_add_pPr(); bdr = OxmlElement("w:pBdr"); bot = OxmlElement("w:bottom")
        bot.set(qn("w:val"), "single"); bot.set(qn("w:sz"), "6"); bot.set(qn("w:space"), "1"); bot.set(qn("w:color"), "D0D7DE")
        bdr.append(bot); ppr.append(bdr)

    def tbl(rows):
        t = doc.add_table(rows=0, cols=len(rows[0])); t.style = "Table Grid"; t.alignment = WD_TABLE_ALIGNMENT.CENTER
        for ri, row in enumerate(rows):
            cells = t.add_row().cells
            for ci, val in enumerate(row):
                cells[ci].text = ""; pp = cells[ci].paragraphs[0]
                if ri == 0:
                    rr = pp.add_run(val); rr.bold = True
                    rr.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF); rr.font.size = Pt(9.5)
                    _lg_shade_cell(cells[ci], "1F6FEB")
                else:
                    _lg_runs(pp, val)
                    for rn in pp.runs: rn.font.size = Pt(9.5)
        doc.add_paragraph()

    # ── Front matter ─────────────────────────────────────────────────────────
    h1(f"{COURSE_NAME} — Step-by-Step Learner Guide")
    p(f"**Course Code:** {COURSE_CODE}  ·  **Version 1.0**  ·  {INSTITUTION}")
    p("Welcome! This guide walks you command-by-command through every hands-on lab. "
      "In one day you deploy your own AI agent platform, connect it to LLMs and messaging "
      "channels, add tools and skills, secure it, and verify financial documents with "
      "blockchain logic — all running on a Hostinger VPS or your local Docker Desktop.")
    note("Work through the labs in order — each builds on the last. Whenever you see a "
         "**Test it** box, stop and confirm the result before moving on.")

    rule()
    h1("0. Before You Start")
    h2("0.1  What you need")
    tbl([
        ["Requirement", "Details", "Where to get it"],
        ["Laptop", "Windows 10/11, macOS 12+, or Ubuntu 22.04+", "Your own machine"],
        ["Docker Desktop (Option B)", "Required for local labs + Lab 10 blockchain", "docker.com/products/docker-desktop"],
        ["Hostinger VPS (Option A)", "Ubuntu 22.04 LTS, SSH access, runs 24/7", REFERRAL],
        ["Telegram account", "Lab 3 — channel setup via BotFather", "telegram.org"],
        ["WhatsApp (phone)", "Lab 3 — QR pairing", "Pre-installed on your phone"],
        ["Groq API key (free)", "Lab 2 — no credit card needed", "console.groq.com"],
        ["AgentMail account (free)", "Lab 4 — email tool", "agentmail.to"],
        ["Firecrawl account (free)", "Lab 4 — web scraper tool", "firecrawl.dev"],
    ])

    h2("0.2  Two ways to run every lab")
    p("**Option A — Hostinger VPS (recommended for 24/7 deployment).**  SSH into your Ubuntu 22.04 "
      "server and run all commands there. OpenClaw stays running after class.")
    p(f"**Hostinger referral (discount):**  {REFERRAL}")
    p("**Option B — Docker Desktop (local laptop).**  Run OpenClaw in a Docker container on your "
      "Windows, macOS, or Linux machine. Requires Docker Desktop to be open and running.")

    h2("0.3  Lab index")
    tbl([["Lab", "Title", "Time"]] +
        [[num, title, f"{mins} min"] for num, title, mins, *_ in LABS])

    # ── Per-lab sections ──────────────────────────────────────────────────────
    for num, title, mins, steps, test_cmd, expected, url in LABS:
        rule()
        h1(f"{num} — {title}")
        p(f"**Estimated time:** {mins} minutes  ·  **Environment:** Hostinger VPS OR Docker Desktop")
        if url:
            p(f"**Lab reference:** {url}")
        h2("Steps")
        for i, (step_label, command) in enumerate(steps, 1):
            p(f"**Step {i} — {step_label}**")
            code(command)
        test_it(f"`{test_cmd}`  →  Expected:  {expected}")

    # ── Reference + Assessment ────────────────────────────────────────────────
    rule()
    h1("Reference Documentation")
    tbl([
        ["Resource", "URL"],
        ["OpenClaw Install",      "https://docs.openclaw.ai/install"],
        ["LLM Providers",         "https://docs.openclaw.ai/providers"],
        ["Channels",              "https://docs.openclaw.ai/channels"],
        ["Tools",                 "https://docs.openclaw.ai/tools"],
        ["Skills / ClawHub",      "https://clawhub.ai"],
        ["Security Guide",        "https://docs.openclaw.ai/gateway/security"],
        ["Dashboard",             "https://docs.openclaw.ai/dashboard"],
        ["Groq (free LLM API)",   "https://console.groq.com"],
        ["Ollama (local model)",  "https://ollama.com"],
        ["AgentMail",             "https://agentmail.to"],
        ["Firecrawl",             "https://www.firecrawl.dev"],
        ["Agent Browser",         "https://agent-browser.dev"],
        ["Hostinger VPS",         REFERRAL],
        ["Course LMS",            "https://ai-lms-tms.tertiaryinfo.tech/"],
    ])

    rule()
    h1("Final Assessment")
    p("**Written Assessment (Short Answer) + Oral Questioning — Open Book.**  "
      "Students must demonstrate the following during the practical assessment:")
    bullets([
        "`openclaw gateway status`  →  Gateway: running",
        "Agent responds to a message via Telegram and/or WhatsApp",
        "Firecrawl tool returns a webpage summary",
        "Cron job fires on schedule and sends result to channel",
        "OpenClaw dashboard accessible at http://localhost:18789",
        "Blockchain: invoice INV-2026-001 registers and verifies via Docker on port 5000",
    ])
    note("TRAQOM Survey and Certificate Delivery are mandatory at end of class. "
         "The trainer will display a QR code — scan it with your phone before leaving.")

    prodoc.add_page_numbers(doc)
    prodoc.enable_update_fields(doc)

    out = os.path.join(COURSEWARE, "LG-WSQ-Business-Transformation-with-OpenClaw.docx")
    doc.save(out)
    print(f"  Saved: {os.path.basename(out)}")
    return out


# ── LP helpers ────────────────────────────────────────────────────────────────
def _lp_shade_cell(cell, hexc):
    tcPr = cell._tc.get_or_add_tcPr(); shd = OxmlElement("w:shd")
    shd.set(qn("w:val"), "clear"); shd.set(qn("w:color"), "auto"); shd.set(qn("w:fill"), hexc); tcPr.append(shd)

def _lp_cell(cell, text, bold=False, color=None, size=10, align=None, italic=False):
    cell.text = ""; pp = cell.paragraphs[0]
    if align is not None: pp.alignment = align
    r = pp.add_run(text); r.bold = bold; r.italic = italic; r.font.size = Pt(size)
    if color: r.font.color.rgb = color

def _lp_schedule_table(doc, rows):
    t = doc.add_table(rows=1, cols=3); t.style = "Table Grid"; t.alignment = WD_TABLE_ALIGNMENT.CENTER
    hdr = t.rows[0].cells
    for i, h in enumerate(["Time", "Topic / Activity", "Duration"]):
        _lp_cell(hdr[i], h, bold=True, color=RGBColor(0xFF,0xFF,0xFF), size=10,
                 align=WD_ALIGN_PARAGRAPH.CENTER if i != 1 else WD_ALIGN_PARAGRAPH.LEFT)
        _lp_shade_cell(hdr[i], "1F6FEB")
    for time_str, activity, minutes, kind in rows:
        c = t.add_row().cells
        _lp_cell(c[0], time_str, size=9.5, align=WD_ALIGN_PARAGRAPH.CENTER, bold=(kind == "topic"))
        _lp_cell(c[1], activity, size=9.5, bold=(kind in ("topic", "break")),
                 color=RGBColor(0x1F,0x6F,0xEB) if kind == "topic" else (RGBColor(0x55,0x5B,0x66) if kind == "break" else None),
                 italic=(kind == "break"))
        _lp_cell(c[2], (f"{minutes} min" if minutes else "—"), size=9.5, align=WD_ALIGN_PARAGRAPH.CENTER)
        fill = "E8F0FE" if kind == "topic" else ("FFF4E5" if kind == "break" else None)
        if fill:
            for cc in c: _lp_shade_cell(cc, fill)
    for row in t.rows:
        for i, w in enumerate([DInches(1.15), DInches(4.5), DInches(0.85)]): row.cells[i].width = w
    doc.add_paragraph()
    return t


# ── LP ───────────────────────────────────────────────────────────────────────
def build_lp():
    ASSETS = os.path.join(REPO, "courseware", "assets")
    VERSIONS = [("1.0", "July 2026",
                 "First version — 1-day lesson plan for WSQ Business Transformation with OpenClaw; "
                 "theory in morning, setup labs before lunch, practice labs in afternoon",
                 INSTITUTION)]

    doc = Document()
    doc.styles["Normal"].font.name = "Arial"; doc.styles["Normal"].font.size = Pt(11)
    prodoc.style_headings(doc)
    prodoc.add_cover_page(doc, "LESSON PLAN", COURSE_NAME, "1.0",
                          org_logo=os.path.join(ASSETS, "tertiary-infotech-logo.png"))
    prodoc.add_version_control(doc, VERSIONS)
    prodoc.add_toc(doc, levels="1-1")

    def heading(text, space_before=12, space_after=6):
        pp = doc.add_paragraph(style="Heading 1")
        pp.paragraph_format.space_before = Pt(space_before); pp.paragraph_format.space_after = Pt(space_after)
        pp.add_run(text)

    def session_label(text):
        pp = doc.add_paragraph(); pp.paragraph_format.space_before = Pt(8); pp.paragraph_format.space_after = Pt(3)
        r = pp.add_run(text); r.bold = True; r.font.size = Pt(11); r.font.color.rgb = RGBColor(0x1F, 0x6F, 0xEB)

    # Course info table
    tbl = doc.add_table(rows=0, cols=2); tbl.style = "Table Grid"; tbl.alignment = WD_TABLE_ALIGNMENT.CENTER
    for k, v in [
        ("Course Title",   COURSE_NAME),
        ("Course Code",    COURSE_CODE),
        ("TGS Ref No",     COURSE_CODE),
        ("Duration",       "1 Day (7 hours)"),
        ("Schedule",       "09:30 AM – 05:00 PM  |  Tea break 10:50  |  Lunch 13:05–14:05"),
        ("Delivery Mode",  "Instructor-led, hands-on labs"),
        ("Environment",    "Option A: Hostinger VPS (Ubuntu 22.04)  OR  Option B: Docker Desktop"),
        ("Institution",    INSTITUTION),
        ("Instructor",     INSTRUCTOR),
    ]:
        row = tbl.add_row().cells
        _lp_cell(row[0], k, bold=True, color=RGBColor(0x11,0x18,0x27), size=10); _lp_shade_cell(row[0], "F1F5FB")
        _lp_cell(row[1], v, size=10)
    for row in tbl.rows:
        row.cells[0].width = DInches(2.1); row.cells[1].width = DInches(4.4)
    doc.add_paragraph()

    heading("Course Overview")
    doc.add_paragraph(
        "This 1-day hands-on WSQ course teaches participants to deploy and operate OpenClaw, "
        "an open-source AI agent framework, on a Hostinger VPS or local Docker Desktop. "
        "The morning covers four theory topics (Overview, Tools & Channels, Security, and "
        "the Blockchain Applications demo) followed by setup labs 1–3 before lunch. "
        "The afternoon moves through practice labs 4–9 covering tools, skills, cron jobs, "
        "security hardening, the web dashboard, and cost optimisation.")

    heading("Learning Outcomes")
    doc.add_paragraph("By the end of this course, participants will be able to:")
    for o in [
        "LO1: Deploy OpenClaw on a Hostinger VPS or Docker Desktop local machine",
        "LO2: Connect LLM providers — Groq (free), OpenAI, Ollama, or Default config",
        "LO3: Set up Telegram (BotFather) and WhatsApp (QR pairing) messaging channels",
        "LO4: Integrate AgentMail, Agent Browser, and Firecrawl tools",
        "LO5: Install and invoke skills from the ClawHub marketplace",
        "LO6: Schedule tasks with cron jobs and monitor uptime with heartbeat",
        "LO7: Apply 10-step security hardening for OpenClaw",
        "LO8: Access the OpenClaw web dashboard locally and via SSH tunnel",
        "LO9: Optimise LLM costs using model selection and context compaction",
        "LO10: Implement blockchain invoice verification using Docker and OpenClaw",
    ]:
        pp = doc.add_paragraph(o, style="List Bullet"); pp.paragraph_format.space_after = Pt(2)

    heading("Prerequisites")
    for pre in [
        "Laptop: Windows 10/11 (WSL2), macOS 12+, or Ubuntu 22.04+",
        "Docker Desktop installed (required for Option B and Lab 10)",
        "Admin / sudo rights on the machine",
        "Telegram account + phone with WhatsApp installed",
        "Free Groq API key at https://console.groq.com (no credit card required)",
        "Free AgentMail account at https://agentmail.to",
        "Free Firecrawl account at https://www.firecrawl.dev",
    ]:
        pp = doc.add_paragraph(pre, style="List Bullet"); pp.paragraph_format.space_after = Pt(2)

    heading("1-Day Schedule", space_before=16)
    session_label("Morning Session  ·  09:30 AM – 01:05 PM")
    _lp_schedule_table(doc, [
        ("09:30 – 09:50", "Welcome, Digital Attendance & Course Overview", 20, "normal"),
        ("09:50 – 10:20", "Topic 1 — Overview of OpenClaw: architecture, use cases & diagrams", 30, "topic"),
        ("10:20 – 10:50", "Topic 3 — Tools, Skills & Channels Overview", 30, "topic"),
        ("10:50 – 11:05", "Tea Break", 15, "break"),
        ("11:05 – 11:25", "Topic 4 — Security, Cost Saving & Best Practices", 20, "topic"),
        ("11:25 – 11:45", "Topic 2 — OpenClaw Applications + Blockchain Demo", 20, "topic"),
        ("11:45 – 12:15", "Lab 1: OpenClaw Hosting (Hostinger VPS or Docker Desktop)", 30, "normal"),
        ("12:15 – 12:45", "Lab 2: OpenClaw Model (Groq free / OpenAI / Ollama)", 30, "normal"),
        ("12:45 – 13:05", "Lab 3: OpenClaw Channel (Telegram + WhatsApp)", 20, "normal"),
    ])
    _lp_schedule_table(doc, [("13:05 – 14:05", "Lunch Break", 60, "break")])
    session_label("Afternoon Session  ·  02:05 PM – 05:00 PM")
    _lp_schedule_table(doc, [
        ("14:05 – 14:35", "Lab 4: OpenClaw Tools (AgentMail / Agent Browser / Firecrawl)", 30, "normal"),
        ("14:35 – 14:55", "Lab 5: OpenClaw Skills (ClawHub marketplace)", 20, "normal"),
        ("14:55 – 15:25", "Lab 6: Cron Jobs and Heartbeat", 30, "normal"),
        ("15:25 – 15:55", "Lab 7: OpenClaw Security (10-step hardening guide)", 30, "normal"),
        ("15:55 – 16:15", "Lab 8: OpenClaw Dashboard (SSH tunnel for VPS users)", 20, "normal"),
        ("16:15 – 16:35", "Lab 9: OpenClaw Cost Saving (model selection + compaction)", 20, "normal"),
        ("16:35 – 17:00", "TRAQOM Survey, Certificate Delivery & Closing", 25, "normal"),
    ])

    heading("Labs Covered", space_before=16)
    doc.add_paragraph("All 10 hands-on labs run on the student's chosen environment (Hostinger VPS or Docker Desktop):")
    rt = doc.add_table(rows=0, cols=3); rt.style = "Table Grid"; rt.alignment = WD_TABLE_ALIGNMENT.CENTER
    hr = rt.add_row().cells
    for i, h in enumerate(["Lab", "Title", "Duration"]):
        _lp_cell(hr[i], h, bold=True, color=RGBColor(0xFF,0xFF,0xFF), size=9.5)
        _lp_shade_cell(hr[i], "1F6FEB")
    for num, title, mins, *_ in LABS:
        c = rt.add_row().cells
        _lp_cell(c[0], num, size=9, align=WD_ALIGN_PARAGRAPH.CENTER)
        _lp_cell(c[1], title, size=9)
        _lp_cell(c[2], f"{mins} min", size=9, align=WD_ALIGN_PARAGRAPH.CENTER)
    for row in rt.rows:
        for i, w in enumerate([DInches(0.8), DInches(4.6), DInches(0.9)]): row.cells[i].width = w

    heading("Assessment", space_before=14)
    doc.add_paragraph(
        "Participants are assessed through hands-on lab work throughout the day and a final "
        "Written Assessment (Short Answer) + Oral Questioning on the LMS (open book). "
        "A minimum of 75% attendance is required for SSG funding eligibility. TRAQOM survey is mandatory.")
    doc.add_paragraph("Practical demonstration — students must show:")
    for a in [
        "openclaw gateway status  →  Gateway: running",
        "Agent responds to a message via Telegram and/or WhatsApp",
        "Firecrawl tool returns a webpage summary",
        "Cron job fires on schedule and sends result to channel",
        "Dashboard accessible at http://localhost:18789",
        "Blockchain: invoice INV-2026-001 registers and verifies via Docker on port 5000",
    ]:
        pp = doc.add_paragraph(a, style="List Bullet"); pp.paragraph_format.space_after = Pt(2)

    prodoc.add_page_numbers(doc)
    prodoc.enable_update_fields(doc)

    out = os.path.join(COURSEWARE, "LP-WSQ-Business-Transformation-with-OpenClaw.docx")
    doc.save(out)
    print(f"  Saved: {os.path.basename(out)}")
    return out


# ── Main ─────────────────────────────────────────────────────────────────────
if __name__ == "__main__":
    print("\n=== Phase 1: Rendering PDF pages ===")
    page_images = render_pdf_pages()

    print("\n=== Building PPT ===")
    build_ppt(page_images)

    print("\n=== Building LG ===")
    lg_path = build_lg()

    print("\n=== Building LP ===")
    lp_path = build_lp()

    print("\n  (PDF conversion skipped — open each .docx in Word and File > Save As > PDF)")

    print("\n=== Done - files in courseware/ ===")
    for f in sorted(os.listdir(COURSEWARE)):
        if not f.startswith("~$"):
            sz = os.path.getsize(os.path.join(COURSEWARE, f)) // 1024
            print(f"  {f}  ({sz:,} KB)")
